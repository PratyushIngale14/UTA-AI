import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

def chunk_text(text, chunk_size=500):
    words = text.split()
    return [" ".join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]

def extract_text_from_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return chunk_text(text)

def extract_text_from_docx(file):
    doc = Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    return chunk_text(text)

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return chunk_text(text)

def load_file(uploaded_file):
    filename = uploaded_file.name.lower()
    if filename.endswith(".pdf"):
        return extract_text_from_pdf(uploaded_file)
    elif filename.endswith(".docx"):
        return extract_text_from_docx(uploaded_file)
    elif filename.endswith(".pptx"):
        return extract_text_from_pptx(uploaded_file)
    else:
        raise ValueError("Unsupported file type.")
